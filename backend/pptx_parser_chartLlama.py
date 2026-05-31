"""PPTX parser with native chart extraction and image analysis pipeline."""

import json
import logging
from PIL import Image
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Optional
import re
from backend.Graph.PipelineMin import extract_chart_with_groq, extract_chart_with_ollama, extract_chart
# from chart_ocr_extractor import describe_image  # ancien pipeline LLaVA/ocr
from backend.Graph.MeuilleurVersionGraph import describe_image_groq

logger = logging.getLogger(__name__)

OLLAMA_URL = "http://localhost:11434/api/chat"

AGENDA_TITLE_HINTS = (
    "ordre du jour",
    "agenda",
    "sommaire",
    "plan",
)

AGENDA_ITEMS_TO_IGNORE = (
)

_pending_images: dict[str, list[dict]] = {}


def _normalize_text(text: Optional[str]) -> str:
    if not text:
        return ""
    normalized = text.lower()
    normalized = normalized.replace("’", "'")
    normalized = " ".join(normalized.split())
    return normalized


def _is_agenda_slide(slide_dict: dict) -> bool:
    title = _normalize_text(slide_dict.get("titre"))
    content_items = slide_dict.get("contenu") or []
    content_joined = " ".join(_normalize_text(item) for item in content_items)
    haystack = f"{title} {content_joined}".strip()
    return any(hint in haystack for hint in AGENDA_TITLE_HINTS)


def _extract_agenda_items(slide_dict: dict) -> list[str]:
    content_items = slide_dict.get("contenu") or []
    cleaned = []
    for item in content_items:
        text = (item or "").strip()
        if not text:
            continue
        low = _normalize_text(text)
        if low in AGENDA_TITLE_HINTS:
            continue
        # Compare against les éléments ignorés après normalisation
        if any(_normalize_text(ignored) in low for ignored in AGENDA_ITEMS_TO_IGNORE):
            continue
        # Ignore numéros de points isolés (ex: "1", "02")
        if re.fullmatch(r"\d{1,3}", low):
            continue
        # Ignore durées isolées (ex: "10 min", "25 mn", "1h30")
        if re.fullmatch(r"\d+\s*(min|mn|h|heure|heures)(\s*\d+\s*min)?", low):
            continue
        # Garde seulement les lignes qui contiennent de vraies lettres
        if not re.search(r"[a-zà-ÿ]", low):
            continue
        cleaned.append(text)
    return cleaned


def _is_transition_slide(slide_dict: dict) -> bool:
    """
    Détecte une slide de transition qui ne contient que le titre d'ordre du jour.
    Retourne True si la slide a un titre correspondant à un hint d'agenda
    et ne contient pas d'autres contenus (texte, graphiques, images, tableaux, notes).
    """
    title = _normalize_text(slide_dict.get("titre"))
    if not title:
        return False

    # Si la slide contient des éléments riches, ce n'est pas une transition
    if slide_dict.get("contenu") or slide_dict.get("graphiques") or slide_dict.get("images") or slide_dict.get("tableaux") or slide_dict.get("notes"):
        # Vérifie qu'il n'y a que des valeurs vides dans contenu
        if any((c or "").strip() for c in (slide_dict.get("contenu") or [])):
            return False

    return any(hint in title for hint in AGENDA_TITLE_HINTS)


def _assign_agenda_to_slides(slides_data: list[dict]) -> list[dict]:
    agenda_slide_idx = None
    agenda_items = []

    # 1. Trouver la slide agenda
    for idx, slide in enumerate(slides_data):
        if _is_agenda_slide(slide):
            agenda_items = _extract_agenda_items(slide)
            if agenda_items:
                agenda_slide_idx = idx
                break

    if not agenda_items or agenda_slide_idx is None:
        return slides_data

    # 2. Slides après agenda — sauter les slides de transition (titre seul)
    start_idx = agenda_slide_idx + 1
    while start_idx < len(slides_data) and _is_transition_slide(slides_data[start_idx]):
        start_idx += 1
    content_slides = slides_data[start_idx:]

    n = len(content_slides)
    k = len(agenda_items)

    if k == 0:
        return slides_data

    # 3. Taille des blocs
    block_size = max(1, n // k)

    current_item_idx = 0

    for i, slide in enumerate(content_slides):
        slide["ordre du jour"] = agenda_items[current_item_idx]

        # passage au point suivant
        if (i + 1) % block_size == 0 and current_item_idx < k - 1:
            current_item_idx += 1

    # 4. Nettoyage slide agenda
    slides_data[agenda_slide_idx]["ordre du jour"] = None

    return slides_data


def _extract_charts(slide) -> list[dict]:
    """
    Extrait les graphiques natifs PowerPoint (éditables).
    Retourne les données brutes depuis le XML — pas de LLM, pas d'OCR.
    Gère : barres, courbes, camembert, aire, radar, scatter, doughnut.
    """
    charts = []
    for shape in slide.shapes:
        if shape.has_chart:
            charts.append(_parse_chart(shape.chart, shape.name))
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            charts.extend(_extract_charts_from_group(shape.shapes))
    return charts


def _extract_charts_from_group(shapes) -> list[dict]:
    charts = []
    for shape in shapes:
        if shape.has_chart:
            charts.append(_parse_chart(shape.chart, shape.name))
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            charts.extend(_extract_charts_from_group(shape.shapes))
    return charts


def _parse_chart(chart, shape_name: str) -> dict:
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type_map = {
        XL_CHART_TYPE.BAR_CLUSTERED:        "barres groupées (horizontal)",
        XL_CHART_TYPE.BAR_STACKED:          "barres empilées (horizontal)",
        XL_CHART_TYPE.BAR_STACKED_100:      "barres empilées 100% (horizontal)",
        XL_CHART_TYPE.COLUMN_CLUSTERED:     "colonnes groupées (vertical)",
        XL_CHART_TYPE.COLUMN_STACKED:       "colonnes empilées",
        XL_CHART_TYPE.COLUMN_STACKED_100:   "colonnes empilées 100%",
        XL_CHART_TYPE.LINE:                 "courbes",
        XL_CHART_TYPE.LINE_MARKERS:         "courbes avec marqueurs",
        XL_CHART_TYPE.LINE_STACKED:         "courbes empilées",
        XL_CHART_TYPE.PIE:                  "camembert",
        XL_CHART_TYPE.PIE_EXPLODED:         "camembert éclaté",
        XL_CHART_TYPE.DOUGHNUT:             "anneau",
        XL_CHART_TYPE.AREA:                 "aire",
        XL_CHART_TYPE.AREA_STACKED:         "aire empilée",
        XL_CHART_TYPE.XY_SCATTER:           "nuage de points",
        XL_CHART_TYPE.XY_SCATTER_LINES:     "nuage de points avec lignes",
        XL_CHART_TYPE.RADAR:                "radar",
        XL_CHART_TYPE.RADAR_FILLED:         "radar rempli",
        XL_CHART_TYPE.BUBBLE:               "bulles",
    }

    try:
        chart_type = chart_type_map.get(chart.chart_type, str(chart.chart_type))
    except Exception:
        chart_type = "inconnu"

    # Titre
    title = None
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip() or None
    except Exception:
        pass

    # Labels axes
    x_axis_label, y_axis_label = None, None
    try:
        if chart.category_axis.has_title:
            x_axis_label = chart.category_axis.axis_title.text_frame.text.strip()
    except Exception:
        pass
    try:
        if chart.value_axis.has_title:
            y_axis_label = chart.value_axis.axis_title.text_frame.text.strip()
    except Exception:
        pass

    # ── Catégories communes — FIX : depuis la 1ère série seulement ────
    categories = _safe_chart_categories(chart)

    # ── Séries ────────────────────────────────────────────────────────
    series_data = []
    try:
        for series in chart.series:
            series_data.append({
                "nom":        _safe_series_name(series),
                "valeurs":    _safe_values(series),
                "categories": categories,   # ← même référence pour toutes
            })
    except Exception as e:
        series_data = [{"erreur": str(e)}]

    summary = _build_chart_summary(title, chart_type, categories, series_data)

    return {
        "source":     "natif",
        "shape_name": shape_name,
        "type":       chart_type,
        "titre":      title,
        "axe_x":      x_axis_label,
        "axe_y":      y_axis_label,
        "categories": categories,          # ← 5 items, plus 15
        "series":     series_data,
        "resume_pv":  summary,
    }


def _safe_series_name(series) -> str:
    try:
        return series.name or "Série sans nom"
    except Exception:
        return "Série sans nom"


def _safe_values(series) -> list:
    try:
        return [
            round(v, 4) if v is not None else None
            for v in series.values
        ]
    except Exception:
        return []


def _safe_categories(series, chart) -> list:
    """Categories depuis la série ou depuis le chart."""
    try:
        cats = list(series.data_labels) if hasattr(series, "data_labels") else []
    except Exception:
        cats = []
    if not cats:
        cats = _safe_chart_categories(chart)
    return cats


def _safe_chart_categories(chart) -> list:
    """
    Extrait les catégories (labels axe X) UNIQUEMENT depuis la première série.
    Évite la duplication qui survenait en concaténant les catégories
    de toutes les séries.
    """
    try:
        from pptx.oxml.ns import qn
        chart_xml = chart._element

        # Cherche les catégories dans la PREMIÈRE série seulement
        # (elles sont identiques pour toutes les séries d'un même graphique)
        first_ser = chart_xml.find(".//" + qn("c:ser"))
        if first_ser is None:
            return []

        # Catégories texte (strRef)
        cat_strs = first_ser.findall(
            qn("c:cat") + "/" + qn("c:strRef") + "/" +
            qn("c:strCache") + "/" + qn("c:pt") + "/" + qn("c:v")
        )
        if cat_strs:
            return [el.text for el in cat_strs if el.text]

        # Catégories numériques (numRef)
        cat_nums = first_ser.findall(
            qn("c:cat") + "/" + qn("c:numRef") + "/" +
            qn("c:numCache") + "/" + qn("c:pt") + "/" + qn("c:v")
        )
        if cat_nums:
            return [el.text for el in cat_nums if el.text]

        # Fallback : cherche dans tout le XML du chart (ancien comportement)
        # mais déduplique en préservant l'ordre
        cat_strs_all = chart_xml.findall(
            ".//" + qn("c:cat") + "/" + qn("c:strRef") + "/" +
            qn("c:strCache") + "/" + qn("c:pt") + "/" + qn("c:v")
        )
        if cat_strs_all:
            seen, unique = set(), []
            for el in cat_strs_all:
                if el.text and el.text not in seen:
                    seen.add(el.text)
                    unique.append(el.text)
            return unique

    except Exception as e:
        logger.debug(f"_safe_chart_categories error: {e}")

    return []

def _build_chart_summary(title: str, chart_type: str,
                         categories: list, series: list) -> str:
    """
    Construit un résumé textuel structuré pour rédiger un PV.
    Format identique à describe_image() pour cohérence.
    """
    lines = [
        f"TYPE DE GRAPHIQUE : {chart_type}",
        f"TITRE : {title or 'non visible'}",
        "",
        "DONNÉES :",
    ]

    if len(series) == 1:
        # Série unique → format simple : label → valeur
        s = series[0]
        vals = s.get("valeurs", [])
        cats = categories or s.get("categories", [])
        for i, v in enumerate(vals):
            label = cats[i] if i < len(cats) else f"Point {i+1}"
            val_str = str(v) if v is not None else "N/A"
            lines.append(f"  {label} → {val_str}")
    else:
        # Plusieurs séries → tableau par série
        for s in series:
            nom = s.get("nom", "?")
            vals = s.get("valeurs", [])
            cats = categories or s.get("categories", [])
            lines.append(f"\n  Série : {nom}")
            for i, v in enumerate(vals):
                label = cats[i] if i < len(cats) else f"Point {i+1}"
                val_str = str(v) if v is not None else "N/A"
                lines.append(f"    {label} → {val_str}")

    # Observations automatiques (min/max)
    try:
        all_vals = [
            v for s in series
            for v in s.get("valeurs", [])
            if v is not None
        ]
        if all_vals and categories:
            max_v = max(all_vals)
            min_v = min(all_vals)
            lines.append(f"\nOBSERVATIONS :")
            lines.append(f"  Valeur max : {max_v}")
            lines.append(f"  Valeur min : {min_v}")
            if len(all_vals) >= 2:
                var = round((all_vals[-1] - all_vals[0]) / all_vals[0] * 100, 1)
                lines.append(f"  Évolution début→fin : {'+' if var >= 0 else ''}{var}%")
    except Exception:
        pass

    return "\n".join(lines)
def _extract_images_fast(slide) -> tuple[list[dict], list[bytes]]:
    """
    Retourne :
    - images_placeholder : liste de dicts avec juste {"status": "pending"}
    - blobs : liste des bytes bruts pour analyse ultérieure
    """
    placeholders = []
    blobs = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blobs.append(shape.image.blob)
            placeholders.append({"status": "pending"})
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            sub_ph, sub_blobs = _extract_images_from_group(shape.shapes)
            placeholders.extend(sub_ph)
            blobs.extend(sub_blobs)

    return placeholders, blobs

def store_pending_images(token: str, images: list[dict]) -> None:
    _pending_images[token] = images

def get_pending_images(token: str) -> list[dict]:
    return _pending_images.pop(token, [])  # pop = auto-nettoyage après lecture


def get_pending_image_blob(token: str, slide_index: int, image_index: int) -> bytes | None:
    for img in _pending_images.get(token, []):
        if img.get("slide_index") == slide_index and img.get("image_index") == image_index:
            return img.get("blob")
    return None


def iter_image_analysis_events(slide_index: int, image_index: int, blob: bytes):
    try:
        # Un seul appel Groq (pas de streaming)
        result = extract_chart_with_groq(blob)

        # Envoie le résultat en un seul chunk
        yield {
            "type": "image_chunk",
            "slide_index": slide_index,
            "image_index": image_index,
            "delta": json.dumps(result, ensure_ascii=False),
        }

        yield {
            "type": "image_done",
            "slide_index": slide_index,
            "image_index": image_index,
            "result": result,
        }

    except Exception as e:
        yield {
            "type": "image_error",
            "slide_index": slide_index,
            "image_index": image_index,
            "error": str(e),
        }

def clear_pending_images(doc_id: str) -> None:
    _pending_images.pop(doc_id, None)

def parse_pptx_fast(file_path: str, doc_id: str) -> dict:
    prs = Presentation(file_path)
    slides_data = []
    all_pending = []

    for idx, slide in enumerate(prs.slides):
        img_placeholders, blobs = _extract_images_fast(slide)

        slide_dict = {
            "index":        idx + 1,
            "titre":        _extract_title(slide),
            "ordre du jour": None,
            "contenu":      _extract_content(slide),
            "tableaux":     _extract_tables(slide),
            "graphiques":   _extract_charts(slide),   # natifs = instantané
            "images":       img_placeholders,          # juste {"status":"pending"}
            "notes":        _extract_notes(slide),
            "est_vide":     False,
        }

        # Stocker les blobs séparément pour le stream
        for img_idx, blob in enumerate(blobs):
            all_pending.append({
                "slide_index": idx + 1,
                "image_index": img_idx,
                "blob":        blob,
            })

        slides_data.append(slide_dict)

    slides_data = _assign_agenda_to_slides(slides_data)

    # Stocker en mémoire pour le stream SSE
    _pending_images[doc_id] = all_pending

    return {
        "nb_slides":            len(slides_data),
        "nb_slides_vides":      sum(1 for s in slides_data if s["est_vide"]),
        "nb_graphiques_natifs": sum(len(s["graphiques"]) for s in slides_data),
        "nb_images_pending":    len(all_pending),
        "slides":               slides_data,
    }

def parse_pptx(file_path: str) -> dict:
    prs = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        charts_natifs = _extract_charts(slide)      
        images_desc   = _extract_images(slide)    

        slide_dict = {
            "index":          idx + 1,
            "titre":          _extract_title(slide),
            "ordre du jour": None,  
            "contenu":        _extract_content(slide),
            "tableaux":       _extract_tables(slide),
            "graphiques":     charts_natifs,         
            "images":         images_desc,
            "notes":          _extract_notes(slide),
            "est_vide":       False,
        }

        has_content = (
            slide_dict["titre"]      or
            slide_dict["contenu"]    or
            slide_dict["tableaux"]   or
            slide_dict["graphiques"] or  
            slide_dict["images"]     or
            slide_dict["notes"]
        )
        slide_dict["est_vide"] = not has_content
        slides_data.append(slide_dict)
    slides_data = _assign_agenda_to_slides(slides_data)
    total_charts  = sum(len(s["graphiques"]) for s in slides_data)
    total_images  = sum(len(s["images"])     for s in slides_data)

    return {
        "nb_slides":        len(slides_data),
        "nb_slides_vides":  sum(1 for s in slides_data if s["est_vide"]),
        "nb_graphiques_natifs": total_charts,    
        "nb_images_ocr":        total_images,
        "slides":               slides_data,
    }



def _extract_images(slide) -> list[dict]:
    images_desc = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            result = _run_chart_pipeline(shape.image.blob)
            images_desc.append(result)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            images_desc.extend(_extract_images_from_group(shape.shapes))
    return images_desc

def _extract_images_from_group(shapes):
    images_desc = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images_desc.append(_run_chart_pipeline(shape.image.blob))
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            images_desc.extend(_extract_images_from_group(shape.shapes))
    return images_desc


def _run_chart_pipeline(image_bytes: bytes) -> dict:
    try:
        #description=extract_chart(image_bytes)
        description = extract_chart_with_ollama(image_bytes)  
        if isinstance(description, dict):
            return description
        return {
            "source": "Qween",
            "description": description,
        }
        #tokens = get_ocr_tokens(image_bytes)
        #result = extract_chart_hybrid(image_bytes, tokens)
        return result or {"erreur": "pipeline returned None"}
    except Exception as e:
        return {"erreur": str(e)}


def _extract_title(slide) -> Optional[str]:
    """
    Cherche le titre dans cet ordre :
    1. Placeholder officiel PowerPoint (type title/center_title)
    2. Première textbox en haut de slide (top < 20% hauteur)
       dont le texte est court (< 120 chars) et la police grande (>= 16pt)
    3. Première textbox en haut si aucune condition de police disponible
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    title_shape = slide.shapes.title
    if title_shape and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()
        if text:
            return text

    slide_height = slide.part.package.presentation_part.presentation.slide_height
    candidates = []

    for shape in slide.shapes:
        if not shape.has_text_frame or shape.has_table or shape.has_chart:
            continue
        if shape.top is None or shape.top > slide_height * 0.25:
            continue

        text = shape.text_frame.text.strip()
        if not text or len(text) > 150:
            continue

        score = shape.top  

        try:
            first_run = shape.text_frame.paragraphs[0].runs
            if first_run and first_run[0].font.size:
                font_pt = first_run[0].font.size.pt
                if font_pt >= 16:
                    score -= slide_height * 0.1  # priorité plus haute
        except Exception:
            pass

        candidates.append((score, text))

    if candidates:
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]

    return None

def _extract_content(slide) -> list[str]:
    """
    Extrait tous les blocs de texte hors titre, tableaux et charts.
    Exclut aussi le texte qui a été retenu comme titre (évite la duplication).
    """
    titre = _extract_title(slide)   # pour l'exclure du contenu
    blocks = []

    shapes_sorted = sorted(
        slide.shapes,
        key=lambda s: (s.top if s.top is not None else 0,
                       s.left if s.left is not None else 0)
    )

    for shape in shapes_sorted:
        if shape == slide.shapes.title:
            continue
        if shape.has_table or shape.has_chart:
            continue
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            # Exclut le texte identique au titre déjà extrait
            if text and text != titre:
                blocks.append(text)

    return blocks


def _extract_tables(slide) -> list[dict]:
    return _extract_tables_from_shapes(slide.shapes)


def _extract_notes(slide) -> Optional[str]:
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf:
            text = notes_tf.text.strip()
            if text and text.lower() not in ("", "click to add notes"):
                return text
    return None


def _extract_tables_from_shapes(shapes):
    tables = []
    for shape in shapes:
        if shape.has_table:
            table = shape.table
            rows_data = [[cell.text.strip() for cell in row.cells]
                         for row in table.rows]
            if rows_data:
                tables.append({
                    "nb_lignes":   len(rows_data),
                    "nb_colonnes": len(rows_data[0]),
                    "lignes":      rows_data,
                })
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            tables.extend(_extract_tables_from_shapes(shape.shapes))
    return tables
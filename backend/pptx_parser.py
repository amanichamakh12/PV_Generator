from pptx import Presentation
from typing import Optional
from pptx.enum.shapes import MSO_SHAPE_TYPE
import cv2
import numpy as np
from PIL import Image
import io
import torch
from ultralytics import YOLO
import os
import re
import platform

# OCR avec pytesseract
try:
    import pytesseract
    import platform
    
    # Configuration pour Windows
    if platform.system() == "Windows":
        # Chemins courants de Tesseract sur Windows
        possible_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Users\stagiaire.risque\AppData\Local\Tesseract-OCR\tesseract.exe"
        ]
        for path in possible_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.pytesseract_cmd = path
                print(f"✅ Tesseract trouvé: {path}")
                break
    
    TESSERACT_AVAILABLE = True
    print("✅ OCR Tesseract disponible")
except ImportError:
    TESSERACT_AVAILABLE = False
    print("⚠️  pytesseract non disponible - OCR désactivé")
except Exception as e:
    TESSERACT_AVAILABLE = False
    print(f"⚠️  Erreur configuration Tesseract: {e}")

def _convert_numpy_types(obj):
    """
    Convertit récursivement les types numpy en types Python pour la sérialisation JSON.
    Gère: integers, floats, booleans, arrays, dicts, listes.
    """
    if isinstance(obj, np.bool_):
        return bool(obj)
    elif isinstance(obj, np.integer):
        return int(obj)
    elif isinstance(obj, np.floating):
        return float(obj)
    elif isinstance(obj, np.ndarray):
        return obj.tolist()
    elif isinstance(obj, dict):
        return {key: _convert_numpy_types(value) for key, value in obj.items()}
    elif isinstance(obj, (list, tuple)):
        return [_convert_numpy_types(item) for item in obj]
    else:
        return obj

# Initialisation des modèles ML
class ChartDetector:
    def __init__(self):
        self.yolo_model = None
        self.cnn_model = None
        self._load_models()

    def _load_models(self):
        """Charge les modèles YOLO et CNN pour la détection de graphiques."""
        try:
            # YOLOv8 pour la détection d'objets (barres, axes, etc.)
            self.yolo_model = YOLO('yolov8n.pt')  # Modèle nano pré-entraîné
            
            # Pour un usage plus avancé, on pourrait entraîner un modèle personnalisé
            # self.yolo_model = YOLO('models/chart_detector.pt')
            
            print("Modèles YOLO chargés avec succès")
        except Exception as e:
            print(f"Erreur lors du chargement des modèles: {e}")
            self.yolo_model = None

# Instance globale du détecteur
chart_detector = ChartDetector()

def parse_pptx(file_path: str) -> dict:
    """
    Parse un fichier .pptx et retourne un dict structuré.
    Gère : texte, tableaux, images, graphes, slides vides, contenu mixte.
    """
    prs = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        slide_dict = {
            "index": idx + 1,
            "titre": _extract_title(slide),
            "contenu": _extract_content(slide),
            "tableaux": _extract_tables(slide),
            "images": _extract_images(slide),
            "graphes": _extract_charts(slide),
            "est_vide": False
        }

        # Marquer les slides sans contenu exploitable
        has_content = (
            slide_dict["titre"] or
            slide_dict["contenu"] or
            slide_dict["tableaux"] or
            slide_dict["images"] or
            slide_dict["graphes"] 

        )
        slide_dict["est_vide"] = not has_content
        slides_data.append(slide_dict)

    return {
        "nb_slides": len(slides_data),
        "nb_slides_vides": sum(1 for s in slides_data if s["est_vide"]),
        "slides": slides_data
    }

def _extract_title(slide) -> Optional[str]:
    """Récupère le titre depuis le placeholder de titre."""
    title_shape = slide.shapes.title
    if title_shape and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()
        return text if text else None
    return None

def _extract_content(slide) -> list[str]:
    """
    Extrait tous les blocs de texte hors titre et hors tableaux.
    Préserve l'ordre visuel (approximatif via top/left).
    """
    blocks = []
    shapes_sorted = sorted(
        slide.shapes,
        key=lambda s: (s.top if s.top is not None else 0,
                       s.left if s.left is not None else 0)
    )

    for shape in shapes_sorted:
        # Ignorer le titre (déjà extrait)
        if shape == slide.shapes.title:
            continue
        # Ignorer les tableaux (traités séparément)
        if shape.has_table:
            continue
        # Ignorer les images (traitées séparément)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        # Ignorer les graphes (traités séparément)
        if shape.has_chart:
            continue
        # Traiter les zones de texte
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    blocks.append(text)

    return blocks

def _extract_tables(slide) -> list[dict]:
    return _extract_tables_from_shapes(slide.shapes)

def _extract_images(slide) -> list[dict]:
    """Extrait les informations sur les images du slide."""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_info = {
                "filename": shape.image.filename if hasattr(shape.image, 'filename') else "unknown",
                "content_type": shape.image.content_type if hasattr(shape.image, 'content_type') else "unknown"
            }
            images.append(image_info)
    return images

def _extract_charts(slide):
    charts = []
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart

            chart_info = {
                "title": chart.chart_title.text_frame.text.strip() if chart.has_title else "",
                "type": str(chart.chart_type),
                "categories": [],
                "series": []
            }

            # Catégories (axe X)
            if chart.plots:
                categories = chart.plots[0].categories
                chart_info["categories"] = [cat.label for cat in categories]

            # Séries (axe Y)
            for series in chart.series:
                chart_info["series"].append({
                    "name": series.name,
                    "values": list(series.values)
                })

            charts.append(chart_info)

    return charts



def _detect_chart_in_image(image_blob: bytes) -> bool:
    """
    Détecte si une image contient un graphique en utilisant YOLO et analyse d'image.
    """
    if not chart_detector.yolo_model:
        # Fallback vers la méthode basique si YOLO n'est pas disponible
        return _detect_chart_basic(image_blob)
    
    try:
        # Convertir les bytes en image PIL
        image = Image.open(io.BytesIO(image_blob))
        
        # Utiliser YOLO pour détecter des objets qui pourraient être des barres
        results = chart_detector.yolo_model(image)
        
        # Analyser les résultats
        detected_objects = []
        for result in results:
            for box in result.boxes:
                class_id = int(box.cls)
                confidence = float(box.conf)
                class_name = result.names[class_id]
                
                # Chercher des objets qui pourraient être des barres ou graphiques
                if confidence > 0.3:
                    detected_objects.append(class_name)
        
        # Si peu d'objets détectés ou objets répétitifs, c'est peut-être un graphique
        unique_objects = set(detected_objects)
        if len(detected_objects) < 5 or len(unique_objects) <= 2:
            # Utiliser aussi l'analyse basique pour confirmer
            return _detect_chart_basic(image_blob)
        
        return False
        
    except Exception as e:
        print(f"Erreur YOLO: {e}")
        return _detect_chart_basic(image_blob)

def _detect_chart_basic(image_blob: bytes) -> bool:
    """
    Méthode basique de détection de graphiques avec OpenCV (fallback).
    """
    try:
        image = Image.open(io.BytesIO(image_blob))
        img_array = np.array(image)
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        edges = cv2.Canny(gray, 50, 150)
        contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        bar_like_contours = []
        for contour in contours:
            x, y, w, h = cv2.boundingRect(contour)
            aspect_ratio = h / w if w > 0 else 0
            if aspect_ratio > 2 and w > 10 and h > 20:
                bar_like_contours.append(contour)
        
        return len(bar_like_contours) >= 3
    except Exception as e:
        print(f"Erreur détection basique: {e}")
        return False


def _extract_chart_data_from_image(image_blob: bytes) -> dict:
    """
    Extrait les données d'un graphique depuis une image en utilisant YOLO et analyse avancée.
    """
    try:
        image = Image.open(io.BytesIO(image_blob))
        img_array = np.array(image)
        
        # Utiliser YOLO pour détecter les barres potentielles
        if chart_detector.yolo_model:
            results = chart_detector.yolo_model(image)
            
            # Pour l'instant, on utilise une approche hybride
            # Idéalement, il faudrait un modèle entraîné spécifiquement pour les graphiques
        
        # Analyse avec OpenCV pour détecter les barres
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        edges = cv2.Canny(gray, 50, 150)
        contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        bars = []
        for contour in contours:
            x, y, w, h = cv2.boundingRect(contour)
            aspect_ratio = h / w if w > 0 else 0
            if aspect_ratio > 1.5 and w > 5 and h > 15:  # Plus permissif pour les barres
                bars.append({
                    "x": int(x),      # Convertir numpy intc en int Python
                    "y": int(y),
                    "width": int(w),
                    "height": int(h),
                    "area": int(w * h)
                })
        
        # Trier par position X (de gauche à droite)
        bars.sort(key=lambda b: b["x"])
        
        # Filtrer les barres qui sont probablement du même graphique
        if bars:
            # Calculer la hauteur moyenne pour normaliser
            avg_height = sum(b["height"] for b in bars) / len(bars)
            max_height = max(b["height"] for b in bars)
            
            # Estimer les valeurs relatives
            values = []
            for bar in bars:
                # Normaliser par rapport à la hauteur max ou moyenne
                relative_value = (bar["height"] / max_height) * 100
                values.append(round(relative_value, 1))
            
            # Essayer de détecter les axes (très basique)
            axes_info = _detect_axes(img_array)
            
            # Reconstruire les données complètes
            reconstructed = _reconstruct_chart_data({
                "detected_bars": len(bars),
                "estimated_values": values,
                "bars_details": bars,
                "axes_detected": axes_info,
                "method": "hybrid_yolo_opencv"
            })
            
            # Analyser les régions texte pour les labels
            text_analysis = _extract_text_regions(image_blob, bars)
            
            # Extraire les données OCR si disponible
            ocr_data = {}
            if TESSERACT_AVAILABLE:
                # Extraire le titre du graphique
                chart_title = _extract_chart_title(image_blob)
                
                # Extraire les labels et valeurs de chaque barre
                bar_labels_values = _extract_bar_labels_and_values(image_blob, bars, text_analysis)
                
                ocr_data = {
                    "ocr_enabled": True,
                    "chart_title": chart_title,
                    "bar_labels_values": bar_labels_values,
                    "method": "pytesseract"
                }
            else:
                ocr_data = {"ocr_enabled": False, "reason": "pytesseract non disponible"}
            
            return _convert_numpy_types({
                **reconstructed,
                "text_regions": text_analysis,
                "ocr_extraction": ocr_data,
                "note": "Données reconstruites avec YOLO + OpenCV + OCR - positions relatives et valeurs extraites"
            })
        else:
            return _convert_numpy_types({
                "detected_bars": 0,
                "estimated_values": [],
                "error": "Aucune barre détectée"
            })
            
    except Exception as e:
        return _convert_numpy_types({"error": str(e)})

def _detect_axes(img_array: np.ndarray) -> dict:
    """
    Détecte les axes du graphique (très basique).
    """
    try:
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        
        # Détection de lignes avec Hough
        edges = cv2.Canny(gray, 50, 150)
        lines = cv2.HoughLinesP(edges, 1, np.pi/180, threshold=100, minLineLength=100, maxLineGap=10)
        
        axes_lines = []
        if lines is not None:
            for line in lines:
                x1, y1, x2, y2 = line[0]
                # Calculer l'angle de la ligne
                angle = np.arctan2(y2 - y1, x2 - x1) * 180 / np.pi
                
                # Lignes horizontales (axe X) ou verticales (axe Y)
                if abs(angle) < 10 or abs(angle - 180) < 10:  # Horizontal
                    axes_lines.append({
                        "type": "horizontal", 
                        "x1": int(x1), 
                        "y1": int(y1), 
                        "x2": int(x2), 
                        "y2": int(y2)
                    })
                elif abs(abs(angle) - 90) < 10:  # Vertical
                    axes_lines.append({
                        "type": "vertical", 
                        "x1": int(x1), 
                        "y1": int(y1), 
                        "x2": int(x2), 
                        "y2": int(y2)
                    })
        
        return _convert_numpy_types({
            "detected_lines": len(axes_lines),
            "axes_lines": axes_lines[:5]  # Limiter à 5 lignes max
        })
        
    except Exception as e:
        return _convert_numpy_types({"error": str(e)})


def _reconstruct_chart_data(chart_data: dict) -> dict:
    """
    Reconstruit les données complètes du graphique à partir des barres détectées.
    Calcule les positions relatives, dimensions et valeurs normalisées.
    """
    if "error" in chart_data or chart_data.get("detected_bars", 0) == 0:
        return chart_data
    
    try:
        bars = chart_data.get("bars_details", [])
        
        if bars:
            # Calculer les positions et dimensions
            y_positions = [bar["y"] for bar in bars]
            heights = [bar["height"] for bar in bars]
            x_positions = [bar["x"] for bar in bars]
            
            y_base = max(y_positions) + min(heights)
            y_top = min(y_positions)
            total_height = y_base - y_top
            max_height = max(heights)
            
            # Reconstruire chaque barre avec contexte
            reconstructed_bars = []
            for idx, bar in enumerate(bars):
                relative_height = (bar["height"] / max_height) * 100 if max_height > 0 else 0
                position_ratio = (bar["x"] - min(x_positions)) / (max(x_positions) - min(x_positions)) if len(bars) > 1 else 0
                
                reconstructed_bars.append({
                    "index": idx,
                    "bar_height_pixels": bar["height"],
                    "bar_width_pixels": bar["width"],
                    "position_x": bar["x"],
                    "position_y": bar["y"],
                    "normalized_value_percent": round(relative_height, 1),
                    "position_ratio": round(position_ratio, 2),
                    "area_pixels": bar["area"]
                })
            
            return {
                "detected_bars": len(bars),
                "bars_reconstruction": reconstructed_bars,
                "chart_dimensions": {
                    "total_height_pixels": total_height,
                    "y_base": int(y_base),
                    "y_top": int(y_top),
                    "max_bar_height": int(max_height)
                },
                "estimated_values": chart_data.get("estimated_values", []),
                "axes_detected": chart_data.get("axes_detected", {}),
                "method": "hybrid_yolo_opencv_reconstruction"
            }
    except Exception as e:
        return {**chart_data, "reconstruction_error": str(e)}


def _extract_text_regions(image_blob: bytes, bars_details: list) -> dict:
    """
    Analyse les régions texte du graphique sans OCR avancé.
    Détecte les zones probables de labels des barres et axes.
    """
    try:
        image = Image.open(io.BytesIO(image_blob))
        img_array = np.array(image)
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        
        # Seuil pour détection de texte
        _, binary = cv2.threshold(gray, 127, 255, cv2.THRESH_BINARY_INV)
        
        height, width = img_array.shape[:2]
        
        # Détection du titre (haut du graphique)
        top_region = binary[:int(height * 0.15), :]
        has_title = bool(np.sum(top_region) > 1000) if top_region.size > 0 else False
        
        # Analyser les régions autour des barres
        bar_text_regions = []
        if bars_details:
            for idx, bar in enumerate(bars_details):
                x, y, w, h = bar["x"], bar["y"], bar["width"], bar["height"]
                
                # Région sous la barre (labels X)
                bottom_start = min(y + h + 5, height - 1)
                bottom_end = min(y + h + 35, height)
                bottom_region = binary[bottom_start:bottom_end, max(0, x - 15):min(x + w + 15, width)]
                
                # Région au-dessus de la barre (valeurs)
                top_start = max(0, y - 25)
                top_end = y
                top_region_bar = binary[top_start:top_end, max(0, x - 15):min(x + w + 15, width)]
                
                text_pixels_below = int(np.sum(bottom_region)) if bottom_region.size > 0 else 0
                text_pixels_top = int(np.sum(top_region_bar)) if top_region_bar.size > 0 else 0
                
                bar_text_regions.append({
                    "bar_index": idx,
                    "position_x": x,
                    "has_label_below": bool(text_pixels_below > 100),
                    "has_value_above": bool(text_pixels_top > 100),
                    "text_pixels_below": text_pixels_below,
                    "text_pixels_above": text_pixels_top
                })
        
        return {
            "detected_title": has_title,
            "bar_text_regions": bar_text_regions,
            "analysis_method": "binary_text_detection"
        }
    except Exception as e:
        return {"text_extraction_error": str(e)}


def _extract_ocr_text(image_blob: bytes, region: dict = None) -> str:
    """
    Extrait le texte d'une image ou d'une région spécifique en utilisant OCR (pytesseract).
    """
    if not TESSERACT_AVAILABLE:
        return ""
    
    try:
        image = Image.open(io.BytesIO(image_blob))
        img_array = np.array(image)
        
        # Si une région est spécifiée, extraire que cette région
        if region:
            x1, y1, x2, y2 = region["x1"], region["y1"], region["x2"], region["y2"]
            roi = img_array[max(0, y1):min(y1 + y2, img_array.shape[0]), 
                           max(0, x1):min(x1 + x2, img_array.shape[1])]
            roi_image = Image.fromarray(roi)
        else:
            roi_image = image
        
        # Configuration OCR pour meilleure reconnaissance
        # --psm 6: Assume a single uniform block of text
        # --oem 3: Use both legacy and LSTM OCR engine modes
        custom_config = r'--psm 6 --oem 3'
        
        # Essayer l'OCR en français + anglais
        text = pytesseract.image_to_string(roi_image, config=custom_config, lang='fra+eng')
        return text.strip()
        
    except Exception as e:
        if "TesseractNotFoundError" in str(type(e)):
            print("⚠️  Tesseract n'est pas installé ou introuvable au chemin configuré")
        else:
            print(f"Erreur OCR: {e}")
        return ""


def _extract_bar_labels_and_values(image_blob: bytes, bars_details: list, text_regions: dict) -> list:
    """
    Extrait les labels et valeurs de chaque barre en utilisant OCR.
    Retourne une liste avec label, valeur textuelle et position pour chaque barre.
    """
    if not TESSERACT_AVAILABLE or not bars_details:
        return []
    
    try:
        image = Image.open(io.BytesIO(image_blob))
        img_array = np.array(image)
        height, width = img_array.shape[:2]
        
        bar_data = []
        
        bar_text_regions = text_regions.get("bar_text_regions", [])
        for region_info in bar_text_regions:
            bar_idx = region_info["bar_index"]
            if bar_idx >= len(bars_details):
                continue
                
            bar = bars_details[bar_idx]
            x, y, w, h = bar["x"], bar["y"], bar["width"], bar["height"]
            
            extracted = {
                "bar_index": bar_idx,
                "position_x": x,
                "label": "",
                "value_above": "",
                "value_pixels": bar["height"]
            }
            
            # Extraire le label sous la barre
            if region_info.get("has_label_below"):
                label_region = {
                    "x1": max(0, x - 20),
                    "y1": min(y + h + 10, height - 1),
                    "x2": min(40, width - x),
                    "y2": min(30, height - (y + h + 10))
                }
                extracted["label"] = _extract_ocr_text(image_blob, label_region)
            
            # Extraire la valeur au-dessus de la barre
            if region_info.get("has_value_above"):
                value_region = {
                    "x1": max(0, x - 20),
                    "y1": max(0, y - 30),
                    "x2": min(w + 40, width - x),
                    "y2": min(35, y)
                }
                value_text = _extract_ocr_text(image_blob, value_region)
                extracted["value_above"] = value_text
                
                # Essayer d'extraire le nombre
                numbers = re.findall(r'\d+', value_text)
                if numbers:
                    extracted["value_numeric"] = int(numbers[0])
            
            bar_data.append(extracted)
        
        return bar_data
        
    except Exception as e:
        print(f"Erreur extraction labels: {e}")
        return []


def _extract_chart_title(image_blob: bytes) -> str:
    """
    Extrait le titre du graphique (région supérieure).
    """
    if not TESSERACT_AVAILABLE:
        return ""
    
    try:
        image = Image.open(io.BytesIO(image_blob))
        img_array = np.array(image)
        height = img_array.shape[0]
        
        # Région du titre (haut du graphique)
        title_region = {
            "x1": 0,
            "y1": 0,
            "x2": img_array.shape[1],
            "y2": min(40, int(height * 0.15))
        }
        
        title_text = _extract_ocr_text(image_blob, title_region)
        return title_text if title_text else ""
        
    except Exception as e:
        print(f"Erreur extraction titre: {e}")
        return ""


# Test de sérialisation JSON pour déboguer les types numpy
if __name__ == "__main__":
    import json
    
    print("🧪 Test de sérialisation JSON avec types numpy...")
    
    # Créer des données de test avec des types numpy
    test_data = {
        "detected_bars": 3,
        "estimated_values": [45.2, 67.8, 23.1],
        "bars_details": [
            {"x": np.int32(10), "y": np.int32(20), "width": np.int32(30), "height": np.int32(40)},
            {"x": np.int32(50), "y": np.int32(20), "width": np.int32(30), "height": np.int32(60)}
        ],
        "axes_detected": {
            "detected_lines": 2,
            "axes_lines": [
                {"type": "horizontal", "x1": np.int32(0), "y1": np.int32(100), "x2": np.int32(200), "y2": np.int32(100)}
            ]
        }
    }
    
    # Convertir et tester la sérialisation
    converted_data = _convert_numpy_types(test_data)
    try:
        json_str = json.dumps(converted_data)
        print("✅ Sérialisation JSON réussie!")
        print(f"Échantillon des données converties: {converted_data['bars_details'][0]}")
    except Exception as e:
        print(f"❌ Erreur de sérialisation: {e}")
        print(f"Types des données: {type(converted_data['bars_details'][0]['x'])}")


def _extract_tables_from_shapes(shapes):
    tables = []

    for shape in shapes:
        # Cas normal
        if shape.has_table:
            table = shape.table
            rows_data = []

            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                # Ignorer les lignes vides (toutes les cellules vides)
                if any(cell for cell in row_cells):
                    rows_data.append(row_cells)

            if rows_data:
                tables.append({
                    "nb_lignes": len(rows_data),
                    "nb_colonnes": len(rows_data[0]),
                    "lignes": rows_data
                })

        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            tables.extend(_extract_tables_from_shapes(shape.shapes))

    return tables

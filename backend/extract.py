"""PPTX text extraction utilities and plain-text PV generation helpers."""

from pptx import Presentation
from datetime import datetime


def extract_text_from_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides, start=1):
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())

        slides_content.append({
            "slide_number": i,
            "content": text_content
        })

    return slides_content


def build_extracted_from_slides(slides_content: list[dict], session_info: dict) -> dict:
    ordre_du_jour = []
    points_discutes = []

    for slide in slides_content:
        title = slide["content"][0] if slide["content"] else f"Slide {slide['slide_number']}"
        ordre_du_jour.append(title)
        points_discutes.append({
            "titre": title,
            "contenu": [str(c) for c in slide["content"][1:]]
        })

    return {
        # ✅ MODIFIÉ : meta vient entièrement de session_info, plus de valeurs par défaut
        "meta": {
            "titre_reunion":    session_info.get("titre_reunion", "Procès-verbal de réunion"),
            "numero":           session_info.get("numero", "PV-2025-001"),
            "date":             session_info.get("date", datetime.now().strftime("%d/%m/%Y")),
            "heure_debut":      session_info.get("heure_debut", ""),
            "heure_fin":        session_info.get("heure_fin", ""),
            "lieu":             session_info.get("lieu", ""),
            "president_seance": session_info.get("president_seance", ""),
            "redacteur":        session_info.get("redacteur", ""),
        },

        # ✅ MODIFIÉ : participants et excusés viennent de session_info
        "participants": session_info.get("participants", []),
        "excuses":      session_info.get("excuses", []),

        # Inchangé — vient toujours de la PPTX
        "ordre_du_jour":  ordre_du_jour,
        "points_discutes": points_discutes,
        "decisions": [],
        "actions":   [],
        "slides":    slides_content
    }

def generate_pv(slides_content: list[dict], session_info: dict) -> str:
    titre = session_info.get("titre_reunion", "Procès-verbal de réunion")
    date  = session_info.get("date", datetime.now().strftime("%d/%m/%Y"))
    lieu  = session_info.get("lieu", "")
    heure_debut = session_info.get("heure_debut", "")
    heure_fin   = session_info.get("heure_fin", "")
    president   = session_info.get("president_seance", "")
    participants = session_info.get("participants", [])

    pv = []
    pv.append(titre)
    pv.append(f"Date : {date}  |  Horaires : {heure_debut} – {heure_fin}  |  Lieu : {lieu}")
    pv.append(f"Présidé par : {president}")
    pv.append("=" * 60 + "\n")

    # Participants
    if participants:
        pv.append("PARTICIPANTS :")
        for p in participants:
            nom = p.get("nom", "") if isinstance(p, dict) else str(p)
            fn  = p.get("fonction", "") if isinstance(p, dict) else ""
            pv.append(f"  - {nom}" + (f" ({fn})" if fn else ""))
        pv.append("")

    pv.append("1. Ordre du jour et contenu des slides\n")
    for slide in slides_content:
        pv.append(f"Slide {slide['slide_number']} :")
        pv.append("\n".join(str(c) for c in slide["content"]))
        pv.append("-" * 40)

    pv.append("\n2. Synthèse de la réunion\n")
    pv.append("Résumé analytique généré par IA.\n")

    return "\n".join(pv)


# Inchangé
def save_pv(pv_text: str, output_file: str = "pv_reunion.txt") -> None:
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(pv_text)